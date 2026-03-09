from flask import Flask, request, send_file
from pdf2docx import Converter
from PIL import Image
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import pymupdf
import mammoth
from weasyprint import HTML
from reportlab.pdfgen import canvas
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
import img2pdf
import os

app = Flask(__name__)

@app.route("/")
def home():
    return "ConvertEase API Running"


# -------------------------
# PDF TO WORD
# -------------------------
@app.route("/pdf-to-word", methods=["POST"])
def pdf_to_word():
    file = request.files["file"]
    file.save("input.pdf")

    cv = Converter("input.pdf")
    cv.convert("output.docx")
    cv.close()

    return send_file("output.docx", as_attachment=True)


# -------------------------
# PDF TO JPG
# -------------------------
@app.route("/pdf-to-jpg", methods=["POST"])
def pdf_to_jpg():
    file = request.files["file"]
    file.save("input.pdf")

    pdf = pymupdf.open("input.pdf")
    page = pdf.load_page(0)
    pix = page.get_pixmap()
    pix.save("output.jpg")

    return send_file("output.jpg", as_attachment=True)


# -------------------------
# JPG TO PDF
# -------------------------
@app.route("/jpg-to-pdf", methods=["POST"])
def jpg_to_pdf():
    file = request.files["file"]
    file.save("input.jpg")

    img = Image.open("input.jpg")
    img.convert("RGB").save("output.pdf")

    return send_file("output.pdf", as_attachment=True)


# -------------------------
# IMAGE TO PDF
# -------------------------
@app.route("/image-to-pdf", methods=["POST"])
def image_to_pdf():
    file = request.files["file"]
    file.save("input_image")

    with open("output.pdf", "wb") as f:
        f.write(img2pdf.convert("input_image"))

    return send_file("output.pdf", as_attachment=True)


# -------------------------
# MERGE PDF
# -------------------------
@app.route("/merge-pdf", methods=["POST"])
def merge_pdf():
    files = request.files.getlist("files")
    merger = PdfMerger()

    for file in files:
        file.save(file.filename)
        merger.append(file.filename)

    merger.write("merged.pdf")
    merger.close()

    return send_file("merged.pdf", as_attachment=True)


# -------------------------
# SPLIT PDF
# -------------------------
@app.route("/split-pdf", methods=["POST"])
def split_pdf():
    file = request.files["file"]
    file.save("input.pdf")

    reader = PdfReader("input.pdf")
    writer = PdfWriter()

    writer.add_page(reader.pages[0])

    with open("split.pdf", "wb") as f:
        writer.write(f)

    return send_file("split.pdf", as_attachment=True)


# -------------------------
# WORD TO PDF
# -------------------------
@app.route("/word-to-pdf", methods=["POST"])
def word_to_pdf():
    file = request.files["file"]
    file.save("input.docx")

    with open("input.docx", "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html = result.value

    with open("temp.html", "w", encoding="utf-8") as f:
        f.write(html)

    HTML("temp.html").write_pdf("output.pdf")

    return send_file("output.pdf", as_attachment=True)


# -------------------------
# PDF COMPRESS
# -------------------------
@app.route("/compress-pdf", methods=["POST"])
def compress_pdf():
    file = request.files["file"]
    file.save("input.pdf")

    reader = PdfReader("input.pdf")
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    with open("compressed.pdf", "wb") as f:
        writer.write(f)

    return send_file("compressed.pdf", as_attachment=True)


# -------------------------
# PDF PROTECT
# -------------------------
@app.route("/protect-pdf", methods=["POST"])
def protect_pdf():
    file = request.files["file"]
    password = request.form.get("password")

    file.save("input.pdf")

    reader = PdfReader("input.pdf")
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(password)

    with open("protected.pdf", "wb") as f:
        writer.write(f)

    return send_file("protected.pdf", as_attachment=True)


# -------------------------
# PDF UNLOCK
# -------------------------
@app.route("/unlock-pdf", methods=["POST"])
def unlock_pdf():
    file = request.files["file"]
    password = request.form.get("password")

    file.save("input.pdf")

    reader = PdfReader("input.pdf")
    reader.decrypt(password)

    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    with open("unlocked.pdf", "wb") as f:
        writer.write(f)

    return send_file("unlocked.pdf", as_attachment=True)


# -------------------------
# EXCEL TO PDF
# -------------------------
@app.route("/excel-to-pdf", methods=["POST"])
def excel_to_pdf():
    file = request.files["file"]
    file.save("input.xlsx")

    df = pd.read_excel("input.xlsx")

    c = canvas.Canvas("excel.pdf")

    y = 800
    for index, row in df.iterrows():
        c.drawString(50, y, str(row.values))
        y -= 20

    c.save()

    return send_file("excel.pdf", as_attachment=True)


# -------------------------
# PPT TO PDF
# -------------------------
@app.route("/ppt-to-pdf", methods=["POST"])
def ppt_to_pdf():
    file = request.files["file"]
    file.save("input.pptx")

    prs = Presentation("input.pptx")

    c = canvas.Canvas("ppt.pdf")
    y = 800

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                c.drawString(50, y, shape.text)
                y -= 20

    c.save()

    return send_file("ppt.pdf", as_attachment=True)


# -------------------------
# RUN SERVER
# -------------------------
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)